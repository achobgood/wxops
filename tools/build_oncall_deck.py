"""Generate the On-Call Connect solution presentation as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
CISCO_BLUE = RGBColor(0x04, 0x9F, 0xD9)
DARK_BG = RGBColor(0x1A, 0x1F, 0x25)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x6B, 0x7B, 0x8D)
GREEN = RGBColor(0x1B, 0xB3, 0x4C)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE8, 0x4D, 0x3D)
DARK_TEXT = RGBColor(0x1A, 0x1F, 0x25)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, text_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color=MID_GRAY):
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

add_text_box(slide, 1, 1.5, 11, 1.2, "On-Call Connect", 48, CISCO_BLUE, True)
add_text_box(slide, 1, 2.8, 11, 0.8, "SMS-Managed On-Call Routing via Webex", 28, WHITE, False)
add_text_box(slide, 1, 4.2, 11, 0.6, "Zero infrastructure. Zero databases. 100% inside Webex.", 20, MID_GRAY)
add_text_box(slide, 1, 5.5, 11, 0.6, "Solution Architecture & Technical Overview", 16, MID_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 12, 0.8, "The Problem", 36, DARK_TEXT, True)

# Problem box
add_rounded_rect(slide, 0.8, 1.6, 11.7, 1.4, RGBColor(0xFE, 0xF3, 0xF2),
                 "", 16, DARK_TEXT)
add_text_box(slide, 1.2, 1.7, 11, 1.2,
             "Customer needs on-call routing for inbound PSTN calls.\n"
             "A dynamic list of cell phone numbers must be dialed sequentially.\n"
             "But they don't want to host, manage, or pay for a database.",
             18, RED)

# What they need
add_text_box(slide, 0.8, 3.4, 12, 0.6, "What they need:", 24, DARK_TEXT, True)

needs = [
    ("Inbound PSTN call", "Callers dial a main number and hear a greeting"),
    ("Dynamic on-call list", "3-5 cell phone numbers that change regularly"),
    ("Sequential dialing", "Try each number for 18 seconds, move to next if no answer"),
    ("Easy updates", "Non-technical staff must be able to change the list"),
    ("No infrastructure", "No servers, no cloud DBs, no external services to manage"),
]
y = 4.2
for title, desc in needs:
    add_rounded_rect(slide, 0.8, y, 0.4, 0.4, CISCO_BLUE)
    add_text_box(slide, 1.5, y - 0.05, 4, 0.5, title, 18, DARK_TEXT, True)
    add_text_box(slide, 5.5, y - 0.05, 7, 0.5, desc, 16, MID_GRAY)
    y += 0.55

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: The Insight
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.5, 12, 0.8, "The Insight", 36, CISCO_BLUE, True)

add_text_box(slide, 0.8, 1.8, 11, 1.5,
             "Webex already has a writable data store:\n"
             "the People API.\n\n"
             "Every user in Control Hub has fields that can be\n"
             "read via GET and updated via PUT.",
             24, WHITE)

# The trick
add_rounded_rect(slide, 0.8, 4.2, 11.7, 2.2, RGBColor(0x0A, 0x3D, 0x5C))
add_text_box(slide, 1.2, 4.3, 11, 0.6, "The Solution", 22, CISCO_BLUE, True)
add_text_box(slide, 1.2, 4.9, 11, 1.4,
             "Create a dedicated service user in Control Hub.\n"
             "Store the on-call phone numbers in the user's title field.\n"
             "Webex Connect reads and writes this field via API.\n"
             "Customer updates the list by sending an SMS.",
             20, WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Solution Architecture (high-level)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 12, 0.8, "Solution Architecture", 36, DARK_TEXT, True)
add_text_box(slide, 0.8, 1.0, 12, 0.5, "Two Webex Connect flows. One service user. Zero external infrastructure.", 18, MID_GRAY)

# Control Hub box (center)
ch_box = add_rounded_rect(slide, 4.8, 2.0, 3.8, 1.8, CISCO_BLUE,
                          "Control Hub\n\nService User\ntitle: +1555...,+1555...", 14, WHITE, True)

# Flow 1 box (left)
f1_box = add_rounded_rect(slide, 0.5, 2.0, 3.5, 1.8, RGBColor(0x2D, 0x8C, 0x5A),
                          "Flow 1: SMS Update\n\nReceive SMS\nParse numbers\nPUT to People API\nConfirm via SMS", 13, WHITE, True)

# Flow 2 box (right)
f2_box = add_rounded_rect(slide, 9.2, 2.0, 3.8, 1.8, RGBColor(0xE8, 0x6E, 0x2C),
                          "Flow 2: Inbound Call\n\nGET People API\nParse title field\nSequential Call Patch\n18s per number", 13, WHITE, True)

# Customer admin (top left)
add_rounded_rect(slide, 0.5, 4.5, 3.5, 0.8, LIGHT_GRAY,
                 "Customer Admin sends SMS", 14, DARK_TEXT, True)

# PSTN Caller (top right)
add_rounded_rect(slide, 9.2, 4.5, 3.8, 0.8, LIGHT_GRAY,
                 "PSTN Caller dials main number", 14, DARK_TEXT, True)

# On-call phones (bottom right)
add_rounded_rect(slide, 9.2, 5.8, 3.8, 0.8, RGBColor(0xE8, 0xF5, 0xE9),
                 "On-Call Cell Phones (sequential dial)", 14, DARK_TEXT, True)

# Arrows (left to center)
add_arrow(slide, 4.0, 2.9, 4.8, 2.9, CISCO_BLUE)
# Arrows (right to center)
add_arrow(slide, 9.2, 2.9, 8.6, 2.9, CISCO_BLUE)
# Admin to flow 1
add_arrow(slide, 2.25, 4.5, 2.25, 3.8, MID_GRAY)
# Caller to flow 2
add_arrow(slide, 11.1, 4.5, 11.1, 3.8, MID_GRAY)
# Flow 2 to on-call
add_arrow(slide, 11.1, 5.8, 11.1, 5.3, GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: SMS Update Flow (detailed)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 12, 0.8, "Flow 1: SMS On-Call Update", 36, DARK_TEXT, True)
add_text_box(slide, 0.8, 1.0, 12, 0.5, "Customer admin texts new on-call numbers. Connect updates Control Hub automatically.", 18, MID_GRAY)

# Step boxes in a horizontal flow
steps = [
    ("1", "Customer Admin\nSends SMS", "5551234567,\n5559876543,\n5557654321", LIGHT_GRAY, DARK_TEXT),
    ("2", "Connect\nReceives SMS", "Receive Node\ncaptures body", RGBColor(0xE3, 0xF2, 0xFD), DARK_TEXT),
    ("3", "Evaluate Node\nValidates", "JavaScript\nformats E.164", RGBColor(0xE8, 0xF5, 0xE9), DARK_TEXT),
    ("4", "HTTP PUT\nPeople API", "Updates title\nfield on user", CISCO_BLUE, WHITE),
    ("5", "SMS Reply\nConfirmation", "On-call updated:\n3 numbers", GREEN, WHITE),
]

x = 0.5
for num, title, desc, bg_color, txt_color in steps:
    add_rounded_rect(slide, x, 1.8, 2.3, 0.6, bg_color, f"Step {num}", 14, txt_color, True)
    add_rounded_rect(slide, x, 2.5, 2.3, 1.0, bg_color, title, 16, txt_color, True)
    add_rounded_rect(slide, x, 3.6, 2.3, 0.9, bg_color, desc, 13, txt_color)
    if x < 10:
        add_arrow(slide, x + 2.3, 3.0, x + 2.6, 3.0, MID_GRAY)
    x += 2.55

# API detail box
add_rounded_rect(slide, 0.5, 5.0, 12.3, 1.8, RGBColor(0xF8, 0xF9, 0xFA))
add_text_box(slide, 0.8, 5.1, 12, 0.5, "API Call Detail", 20, DARK_TEXT, True)
add_text_box(slide, 0.8, 5.6, 12, 1.1,
             'PUT https://webexapis.com/v1/people/{serviceUserId}\n'
             'Authorization: Bearer {service_app_token}\n'
             'Body: {"title": "+15551234567,+15559876543,+15557654321"}',
             14, MID_GRAY, font_name="Consolas")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Inbound Call Flow (detailed)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 12, 0.8, "Flow 2: Inbound Call Routing", 36, DARK_TEXT, True)
add_text_box(slide, 0.8, 1.0, 12, 0.5, "PSTN call arrives. Connect reads on-call list from Control Hub. Dials sequentially.", 18, MID_GRAY)

# Flow steps - two rows
row1 = [
    ("PSTN Call\nArrives", LIGHT_GRAY, DARK_TEXT),
    ("TTS Greeting\n\"Thanks for calling,\nplease wait...\"", RGBColor(0xE3, 0xF2, 0xFD), DARK_TEXT),
    ("GET People API\nRead title field", CISCO_BLUE, WHITE),
    ("Parse CSV\nBuild number array", RGBColor(0xE8, 0xF5, 0xE9), DARK_TEXT),
]

x = 0.5
for text, bg_color, txt_color in row1:
    add_rounded_rect(slide, x, 1.8, 2.8, 1.3, bg_color, text, 15, txt_color, True)
    if x < 9:
        add_arrow(slide, x + 2.8, 2.45, x + 3.1, 2.45, MID_GRAY)
    x += 3.1

# Sequential dial section
add_text_box(slide, 0.5, 3.4, 12, 0.5, "Sequential Dial (18 seconds each):", 20, DARK_TEXT, True)

dial_steps = [
    ("Call Patch #1\n+15551234567\n18s timeout", ORANGE, WHITE),
    ("Call Patch #2\n+15559876543\n18s timeout", ORANGE, WHITE),
    ("Call Patch #3\n+15557654321\n18s timeout", ORANGE, WHITE),
]

x = 0.5
for text, bg_color, txt_color in dial_steps:
    add_rounded_rect(slide, x, 4.0, 3.3, 1.2, bg_color, text, 15, txt_color, True)
    if x < 7:
        # No answer arrow
        add_text_box(slide, x + 3.3, 4.3, 0.8, 0.4, "No\nAnswer", 11, RED)
        add_arrow(slide, x + 3.3, 4.6, x + 4.1, 4.6, RED)
    x += 4.1

# Outcomes
add_rounded_rect(slide, 0.5, 5.7, 5.5, 1.0, GREEN,
                 "Any number answers → Call Connected", 18, WHITE, True)
add_rounded_rect(slide, 6.5, 5.7, 6.3, 1.0, RED,
                 "All numbers fail → TTS: \"Sorry, no one available.\nPlease try again later.\"", 16, WHITE, True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Why This Works
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.5, 12, 0.8, "Why This Works", 36, CISCO_BLUE, True)

benefits = [
    ("No Database Required",
     "Phone numbers are stored in a Webex user's profile field — "
     "the same Control Hub the customer already uses."),
    ("SMS-Driven Updates",
     "Customer admin texts new on-call numbers to a Webex Connect SMS line. "
     "No portal login, no technical skills needed."),
    ("100% Inside Webex",
     "Both flows run in Webex Connect. Data lives in Control Hub. "
     "Authentication uses a Webex service app. Single vendor."),
    ("Real-Time Changes",
     "Text new numbers, they're live on the next call. "
     "No deployment, no sync delay, no cache to clear."),
    ("Auditable",
     "Every update goes through the People API — visible in admin audit logs. "
     "SMS history provides a paper trail of who changed what."),
]

y = 1.6
for title, desc in benefits:
    add_rounded_rect(slide, 0.8, y, 0.15, 0.6, CISCO_BLUE)
    add_text_box(slide, 1.3, y, 4, 0.6, title, 22, WHITE, True)
    add_text_box(slide, 5.5, y, 7.5, 0.7, desc, 16, MID_GRAY)
    y += 1.05

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Technical Requirements
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 12, 0.8, "Technical Requirements", 36, DARK_TEXT, True)

# Requirements table as boxes
headers = ["Component", "Requirement", "Detail"]
header_y = 1.4
add_rounded_rect(slide, 0.5, header_y, 3.0, 0.6, CISCO_BLUE, "Component", 16, WHITE, True)
add_rounded_rect(slide, 3.5, header_y, 4.0, 0.6, CISCO_BLUE, "Requirement", 16, WHITE, True)
add_rounded_rect(slide, 7.5, header_y, 5.3, 0.6, CISCO_BLUE, "Detail", 16, WHITE, True)

rows = [
    ("Service User", "Webex user (free or licensed)", "Dedicated user e.g. oncall-mainline@customer.com"),
    ("Service App", "OAuth with People scopes", "spark-admin:people_read + spark-admin:people_write"),
    ("Webex Connect", "SMS-enabled number", "Receives update texts from customer admin"),
    ("Webex Connect", "Voice inbound number", "PSTN number callers dial"),
    ("Connect Flows", "2 flows (SMS + Voice)", "Built in Webex Connect flow builder"),
    ("Auth Token", "Service App token in Connect", "Stored as flow credential or variable"),
]

y = 2.1
for comp, req, detail in rows:
    bg = LIGHT_GRAY if rows.index((comp, req, detail)) % 2 == 0 else WHITE
    add_rounded_rect(slide, 0.5, y, 3.0, 0.6, bg, comp, 14, DARK_TEXT, True)
    add_rounded_rect(slide, 3.5, y, 4.0, 0.6, bg, req, 14, DARK_TEXT)
    add_rounded_rect(slide, 7.5, y, 5.3, 0.6, bg, detail, 12, MID_GRAY)
    y += 0.6

# Note
add_text_box(slide, 0.5, 5.8, 12.3, 1.0,
             "Note: The service user does not need a Webex Calling license. "
             "It only needs to exist in Control Hub with a writable profile. "
             "The title field is used as a lightweight data store for the on-call number list.",
             16, MID_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: SMS Format & Validation
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 12, 0.8, "SMS Update Interface", 36, DARK_TEXT, True)
add_text_box(slide, 0.8, 1.0, 12, 0.5, "How the customer admin updates the on-call list", 18, MID_GRAY)

# Example conversation
msgs = [
    ("Admin", "5551234567, 5559876543, 5557654321", RGBColor(0xE3, 0xF2, 0xFD), DARK_TEXT),
    ("System", "On-call list updated: 3 numbers\n1. +15551234567\n2. +15559876543\n3. +15557654321", RGBColor(0xE8, 0xF5, 0xE9), DARK_TEXT),
    ("Admin", "5551234567, 5559876543", RGBColor(0xE3, 0xF2, 0xFD), DARK_TEXT),
    ("System", "On-call list updated: 2 numbers\n1. +15551234567\n2. +15559876543", RGBColor(0xE8, 0xF5, 0xE9), DARK_TEXT),
    ("Admin", "hello", RGBColor(0xE3, 0xF2, 0xFD), DARK_TEXT),
    ("System", "Sorry, I didn't recognize any phone numbers.\nPlease send comma-separated numbers.\nExample: 5551234567, 5559876543", RGBColor(0xFE, 0xF3, 0xF2), RED),
]

y = 1.8
for sender, text, bg_color, txt_color in msgs:
    x_offset = 1.5 if sender == "Admin" else 5.0
    width = 6.5
    lines = text.count('\n') + 1
    height = max(0.5, lines * 0.35)
    add_rounded_rect(slide, x_offset, y, width, height, bg_color, text, 14, txt_color)
    add_text_box(slide, x_offset - 0.9 if sender == "Admin" else x_offset + width + 0.1,
                 y, 0.8, 0.4, sender, 12, MID_GRAY, True)
    y += height + 0.15

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Next Steps
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.5, 12, 0.8, "Next Steps", 36, CISCO_BLUE, True)

steps = [
    ("1", "Create Service User", "Create oncall-mainline@customer.com in Control Hub"),
    ("2", "Create Service App", "Register app with People API read/write scopes"),
    ("3", "Build SMS Update Flow", "Webex Connect flow: receive SMS, validate, PUT to API"),
    ("4", "Build Inbound Call Flow", "Webex Connect flow: TTS, GET API, sequential Call Patch"),
    ("5", "Test End-to-End", "Update via SMS, verify in Control Hub, place test call"),
    ("6", "Hand Off to Customer", "Document the SMS number and expected format"),
]

y = 1.6
for num, title, desc in steps:
    add_rounded_rect(slide, 0.8, y, 0.7, 0.7, CISCO_BLUE, num, 24, WHITE, True)
    add_text_box(slide, 1.8, y + 0.05, 4, 0.6, title, 22, WHITE, True)
    add_text_box(slide, 5.8, y + 0.05, 7, 0.6, desc, 18, MID_GRAY)
    y += 0.9

out_path = "/Users/ahobgood/Documents/webexCalling/docs/plans/oncall-connect-solution.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
